#!/usr/bin/env python3
"""
SmartSearch - Ein kleines Tool, das lokale Dateien "versteht"
und dir erlaubt, in normaler Sprache danach zu suchen.

Backend-Modul: search.py (v3.7 - Robustheits-Fixes)

Änderungen gegenüber v3.6:
- OCR verarbeitet jetzt bis zu OCR_MAX_SEITEN Seiten (statt fest 5) - lange
  gescannte Dokumente werden nicht mehr nach Seite 5 stillschweigend
  abgeschnitten.
- geladenes_modell() ist jetzt Thread-sicher (Lock), falls Suche und
  Indexierung gleichzeitig zum allerersten Mal das Modell laden wollen.
- Dateien, die nicht gelesen werden konnten (kaputtes PDF, leerer Inhalt
  etc.) werden jetzt als "verarbeitet, aber ohne Inhalt" markiert, damit sie
  nicht bei jedem Indexierungslauf erneut (erfolglos) verarbeitet werden.
  Sie tauchen weiterhin nicht in Suchergebnissen auf.
"""

import sys
import os
import json
import pickle
import time
import datetime
import subprocess
import re
import threading
import math
from concurrent.futures import ThreadPoolExecutor

INDEX_FILE = os.path.join(os.path.dirname(__file__), "index.pkl")
CONFIG_FILE = os.path.join(os.path.dirname(__file__), "config.json")
FAVORITEN_FILE = os.path.join(os.path.dirname(__file__), "favoriten.json")

UNTERSTUETZT = (".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx")

STOPWOERTER = {
    "der", "die", "das", "den", "dem", "des", "ein", "eine", "einer", "eines",
    "und", "oder", "für", "von", "mit", "auf", "im", "in", "zu", "zur", "zum",
    "ist", "sind", "war", "waren", "als", "am", "an", "bei", "aus", "nach",
}

MODELL_NAME = "BAAI/bge-m3"

# Score-Schwelle für suche_intern(): Treffer unterhalb dieses kombinierten
# Semantik+Keyword-Scores werden verworfen. Empirisch ermittelt (v3.x):
# alles darunter waren in Tests fast immer thematisch irrelevante Treffer.
SCORE_SCHWELLE = 0.15

# Maximale Seitenzahl, die bei einem gescannten PDF per OCR gelesen wird.
# Vorher fest auf 5 begrenzt - dadurch wurden längere Scans (Verträge,
# Berichte) nach Seite 5 stillschweigend nicht mehr durchsucht.
OCR_MAX_SEITEN = 30

# Auflösung (DPI), mit der PDF-Seiten fürs OCR in Bilder umgewandelt
# werden. pdf2image nutzt standardmäßig 200 DPI - für schwierige Scans
# (kleine Schrift, Ausweiskopien, schlecht kontrastierte Dokumente) liefert
# eine höhere Auflösung Tesseract deutlich mehr erkennbare Details. Höher
# = bessere Erkennung, aber langsamer und mehr Arbeitsspeicher pro Seite.
OCR_DPI = 400

# Ordnernamen, die NIE mitindexiert werden sollen, egal wo sie im
# durchsuchten Ordnerbaum auftauchen. Das sind App-eigene/Bibliotheks-
# Ordner (venv, Build-Ausgaben, Python-Cache, Git-Interna) - deren
# enthaltene Dateien (z.B. Vorlagen-.docx/.pptx aus installierten
# Bibliotheken) haben mit den eigentlichen Nutzer-Dokumenten nichts zu
# tun, blähen aber den Index unnötig auf und kosten Rechenzeit.
IGNORIERTE_ORDNERNAMEN = {"venv", ".venv", ".git", "__pycache__", "build", "dist", "node_modules"}

# Anzahl paralleler Threads beim Einlesen der Dateien (Text-Extraktion/OCR).
# Diese Arbeit ist größtenteils I/O bzw. läuft in C-Bibliotheken
# (pdfplumber, Tesseract), die während der Arbeit die Python-GIL freigeben -
# paralleles Lesen bringt hier also einen echten Geschwindigkeitsgewinn,
# unabhängig vom CPU-Modus des KI-Modells (siehe lade_modell()).
LESE_THREADS = 4


def lade_config():
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, "r") as f:
            return json.load(f)
    return {"ordner": []}


def speichere_config(config):
    with open(CONFIG_FILE, "w") as f:
        json.dump(config, f, indent=2)


def befehl_ordner_hinzufuegen(ordner):
    ordner = os.path.abspath(os.path.expanduser(ordner))
    if not os.path.isdir(ordner):
        print(f"Der Ordner '{ordner}' existiert nicht.")
        return
    config = lade_config()
    if ordner in config["ordner"]:
        return
    config["ordner"].append(ordner)
    speichere_config(config)


def befehl_ordner_entfernen(ordner):
    ordner = os.path.abspath(os.path.expanduser(ordner))
    config = lade_config()
    if ordner not in config["ordner"]:
        return
    config["ordner"].remove(ordner)
    speichere_config(config)


# ---------- FAVORITEN ----------

def lade_favoriten():
    """Gibt die Menge der als Favorit markierten Dateipfade zurück."""
    if os.path.exists(FAVORITEN_FILE):
        try:
            with open(FAVORITEN_FILE, "r") as f:
                return set(json.load(f))
        except Exception as e:
            print(f"[Warnung] Favoriten konnten nicht geladen werden: {e}")
            return set()
    return set()


def speichere_favoriten(favoriten):
    try:
        with open(FAVORITEN_FILE, "w") as f:
            json.dump(sorted(favoriten), f, indent=2)
    except Exception as e:
        print(f"[Warnung] Favoriten konnten nicht gespeichert werden: {e}")


def favorit_umschalten(pfad):
    """Fügt pfad zu den Favoriten hinzu oder entfernt ihn.

    Gibt True zurück, wenn die Datei danach ein Favorit ist, sonst False.
    """
    favoriten = lade_favoriten()
    if pfad in favoriten:
        favoriten.remove(pfad)
        ist_favorit = False
    else:
        favoriten.add(pfad)
        ist_favorit = True
    speichere_favoriten(favoriten)
    return ist_favorit


def lade_modell():
    """Lädt das BGE-M3 KI-Modell.

    WICHTIG: Läuft bewusst auf der CPU statt auf MPS (Apples GPU-Backend).
    PyTorchs MPS-Speicherverwalter hat einen bekannten Bug, der bei
    längeren Indexierungsläufen mit "buffer_block INTERNAL ASSERT FAILED"
    abstürzt (die ganze App wird dann vom Betriebssystem beendet - "zsh:
    abort"). Auf der CPU ist die Berechnung etwas langsamer, aber stabil.
    """
    print("Lade KI-Modell (BGE-M3, CPU-Modus)...")
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer(MODELL_NAME, device="cpu")


def _bild_fuer_ocr_aufbereiten(img):
    """Bereitet ein Seitenbild vor dem OCR-Durchlauf auf, um Tesseracts
    Erkennungsrate bei schwierigen Scans (kleine Schrift, wenig Kontrast,
    z.B. Ausweiskopien) zu verbessern.

    - Graustufen: Farbinformation lenkt Tesseract eher ab, als dass sie
      hilft, und macht die Verarbeitung zusätzlich langsamer.
    - Kontrastverstärkung: hebt blasse/schwach gedruckte Zeichen deutlicher
      vom Hintergrund ab.
    """
    from PIL import ImageEnhance
    graustufen = img.convert("L")
    kontrastverstaerkt = ImageEnhance.Contrast(graustufen).enhance(1.6)
    return kontrastverstaerkt


def _pdf_ocr_verarbeiten(pfad):
    """Liest Text von stummen PDFs/Scans per Tesseract-OCR aus.

    Verarbeitet bis zu OCR_MAX_SEITEN Seiten mit OCR_DPI Auflösung, damit
    auch längere gescannte Dokumente und schwierige Scans (kleine Schrift,
    schwacher Kontrast) zuverlässig lesbar sind.
    """
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(pfad, first_page=1, last_page=OCR_MAX_SEITEN, dpi=OCR_DPI)
        ocr_text_teile = []
        for img in images:
            aufbereitet = _bild_fuer_ocr_aufbereiten(img)
            text = pytesseract.image_to_string(aufbereitet, lang="deu+eng")
            if text.strip():
                ocr_text_teile.append(text)
        return "\n".join(ocr_text_teile)
    except Exception as e:
        print(f"[Warnung] OCR fehlgeschlagen bei {pfad}: {e}")
        return ""


def lies_datei(pfad):
    ext = os.path.splitext(pfad)[1].lower()
    try:
        if ext in [".txt", ".md"]:
            with open(pfad, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif ext == ".pdf":
            text = ""
            try:
                import pdfplumber
                with pdfplumber.open(pfad) as pdf:
                    seiten_texte = []
                    for page in pdf.pages:
                        extracted = page.extract_text()
                        if extracted:
                            seiten_texte.append(extracted)
                    text = "\n".join(seiten_texte)
            except Exception as e:
                print(f"[Warnung] pdfplumber fehlgeschlagen bei {pfad}: {e}")

            if not text or len(text.strip()) < 50:
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(pfad)
                    text = "\n".join(page.extract_text() or "" for page in reader.pages)
                except Exception as e:
                    print(f"[Warnung] pypdf-Fallback fehlgeschlagen bei {pfad}: {e}")

            if not text or len(text.strip()) < 50:
                text = _pdf_ocr_verarbeiten(pfad)

            return text
        elif ext == ".docx":
            from docx import Document
            doc = Document(pfad)
            text_teile = [absatz.text for absatz in doc.paragraphs if absatz.text.strip()]

            # WICHTIG: Viele moderne Vorlagen (z.B. Lebenslauf-Layouts mit
            # Spalten für Datum/Firma/Beschreibung) speichern ihren Text in
            # TABELLEN statt in normalen Absätzen. doc.paragraphs erfasst
            # das nicht - solche Dokumente wurden bisher fälschlich als
            # "ohne Inhalt" markiert, obwohl sie voller Text waren.
            for tabelle in doc.tables:
                for zeile in tabelle.rows:
                    for zelle in zeile.cells:
                        for absatz in zelle.paragraphs:
                            if absatz.text.strip():
                                text_teile.append(absatz.text)

            return "\n".join(text_teile)
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(pfad, read_only=True, data_only=True)
            text_teile = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    zeilen_text = " ".join(str(cell) for cell in row if cell is not None)
                    if zeilen_text.strip():
                        text_teile.append(zeilen_text)
            return "\n".join(text_teile)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(pfad)
            text_teile = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_teile.append(shape.text)
            return "\n".join(text_teile)
        else:
            return None
    except Exception as e:
        print(f"[Warnung] Datei konnte nicht gelesen werden: {pfad} ({e})")
        return None


def in_abschnitte_teilen(text, groesse=700, ueberlappung=150):
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=groesse,
            chunk_overlap=ueberlappung,
            separators=["\n\n", "\n", ".", " ", ""]
        )
        return splitter.split_text(text)
    except ImportError:
        woerter = text.split()
        if len(woerter) <= 100:
            return [" ".join(woerter)]
        abschnitte = []
        schritt = max(100 - 20, 1)
        for i in range(0, len(woerter), schritt):
            abschnitte.append(" ".join(woerter[i:i + 100]))
        return abschnitte


def lade_bestehenden_index():
    if os.path.exists(INDEX_FILE):
        with open(INDEX_FILE, "rb") as f:
            return pickle.load(f)
    return []


def speichere_index(eintraege):
    """Speichert den Index atomar: erst in eine temporäre Datei schreiben,
    dann per os.replace() an die Stelle der echten Datei verschieben.

    WICHTIG für die "Suche während der Indexierung"-Funktion: Ohne das
    könnte eine parallel laufende Suche exakt in dem Moment lesen, in dem
    hier gerade geschrieben wird, und eine unvollständige/kaputte Datei
    erwischen. os.replace() ist auf demselben Dateisystem atomar - Leser
    sehen immer entweder die alte oder die neue vollständige Datei, nie
    einen Zwischenzustand.
    """
    temp_pfad = INDEX_FILE + ".tmp"
    with open(temp_pfad, "wb") as f:
        pickle.dump(eintraege, f)
    os.replace(temp_pfad, INDEX_FILE)


def dateien_im_ordner(ordner):
    gefunden = []
    for wurzel, unterordner, dateien in os.walk(ordner):
        # Ignorierte Unterordner direkt aus der Traversierung entfernen
        # (verändert die Liste in-place, damit os.walk gar nicht erst
        # hineinschaut - schneller als hinterher zu filtern).
        unterordner[:] = [u for u in unterordner if u not in IGNORIERTE_ORDNERNAMEN]
        for name in dateien:
            if name.lower().endswith(UNTERSTUETZT):
                gefunden.append(os.path.join(wurzel, name))
    return gefunden


def fehlgeschlagene_dateien():
    """Gibt die Liste aller Dateipfade zurück, die beim letzten
    Indexierungslauf nicht gelesen werden konnten (ohne_inhalt-Flag).

    Für die "⚠️ X Dateien nicht lesbar"-Anzeige in der GUI, damit sowas
    nicht mehr stillschweigend im Terminal verschwindet.
    """
    eintraege = lade_bestehenden_index()
    gesehen = set()
    ergebnis = []
    for e in eintraege:
        if e.get("ohne_inhalt") and e["datei"] not in gesehen:
            gesehen.add(e["datei"])
            ergebnis.append(e["datei"])
    return sorted(ergebnis)


def entferne_fehlgeschlagene_markierung():
    """Entfernt alle 'ohne_inhalt'-Einträge komplett aus dem Index.

    Danach werden diese Dateien beim nächsten aktualisiere_index()-Lauf
    ganz normal erneut versucht (z.B. sinnvoll, nachdem man OCR
    nachinstalliert hat). Gibt die Anzahl der entfernten Einträge zurück.
    """
    eintraege = lade_bestehenden_index()
    verbleibend = [e for e in eintraege if not e.get("ohne_inhalt")]
    entfernt = len(eintraege) - len(verbleibend)
    if entfernt:
        speichere_index(verbleibend)
    return entfernt


def exportiere_konfiguration(zielpfad):
    """Schreibt überwachte Ordner + Favoriten als lesbare JSON-Datei.

    Praktisch als Backup oder um die eigene Einrichtung auf einen anderen
    Mac zu übertragen, ohne den kompletten (oft sehr großen) Suchindex
    mitnehmen zu müssen - der wird beim nächsten 'Index aktualisieren'
    einfach neu aufgebaut.
    """
    daten = {
        "ordner": lade_config().get("ordner", []),
        "favoriten": sorted(lade_favoriten()),
    }
    with open(zielpfad, "w", encoding="utf-8") as f:
        json.dump(daten, f, indent=2, ensure_ascii=False)


def importiere_konfiguration(quellpfad):
    """Liest eine mit exportiere_konfiguration() erzeugte JSON-Datei ein
    und fügt überwachte Ordner + Favoriten zur bestehenden Einrichtung
    hinzu (überschreibt nichts, ergänzt nur). Gibt (anzahl_ordner,
    anzahl_favoriten) aus der importierten Datei zurück."""
    with open(quellpfad, "r", encoding="utf-8") as f:
        daten = json.load(f)

    importierte_ordner = daten.get("ordner", [])
    config = lade_config()
    for ordner in importierte_ordner:
        if ordner not in config["ordner"]:
            config["ordner"].append(ordner)
    speichere_config(config)

    importierte_favoriten = daten.get("favoriten", [])
    favoriten = lade_favoriten()
    favoriten.update(importierte_favoriten)
    speichere_favoriten(favoriten)

    return len(importierte_ordner), len(importierte_favoriten)


def aktualisiere_index(ordner, modell=None, still=False, fortschritt_fn=None):
    ordner = os.path.abspath(os.path.expanduser(ordner))
    alle_eintraege = lade_bestehenden_index()

    eintraege_dieser_ordner = [e for e in alle_eintraege if e["datei"].startswith(ordner + os.sep)]
    andere_eintraege = [e for e in alle_eintraege if not e["datei"].startswith(ordner + os.sep)]

    bekannt = {e["datei"]: e.get("geaendert", 0) for e in eintraege_dieser_ordner}
    aktuelle_dateien = set(dateien_im_ordner(ordner))
    zu_verarbeiten = [p for p in aktuelle_dateien if p not in bekannt or os.path.getmtime(p) > bekannt[p]]

    eintraege_dieser_ordner = [
        e for e in eintraege_dieser_ordner
        if e["datei"] in aktuelle_dateien and e["datei"] not in zu_verarbeiten
    ]

    if not zu_verarbeiten:
        ergebnis = andere_eintraege + eintraege_dieser_ordner
        speichere_index(ergebnis)
        return ergebnis

    if modell is None:
        modell = geladenes_modell()

    neue_eintraege = []

    def _datei_verarbeiten(pfad):
        geaendert = os.path.getmtime(pfad)
        text = lies_datei(pfad)
        return pfad, geaendert, text

    fertig_zaehler = 0
    with ThreadPoolExecutor(max_workers=LESE_THREADS) as pool:
        for pfad, geaendert, text in pool.map(_datei_verarbeiten, zu_verarbeiten):
            fertig_zaehler += 1
            if fortschritt_fn:
                fortschritt_fn(fertig_zaehler, os.path.basename(pfad))

            if not text or not text.strip():
                # Datei konnte nicht gelesen werden oder ist leer (z.B.
                # kaputtes PDF, gescanntes Dokument ohne erkennbaren Text).
                # Trotzdem als "verarbeitet" markieren (mit geaendert-
                # Zeitstempel, aber ohne Vektor), damit sie beim nächsten
                # Lauf nicht erneut - erfolglos - verarbeitet wird. In der
                # Suche taucht sie wegen des fehlenden Vektors nicht auf
                # (siehe suche_intern-Filter). Über fehlgeschlagene_dateien()
                # kann man diese Liste einsehen, über
                # entferne_fehlgeschlagene_markierung() einen Neuversuch
                # erzwingen (z.B. nach nachträglicher OCR-Installation).
                neue_eintraege.append({
                    "datei": pfad,
                    "text": "",
                    "text_fuer_analyse": None,
                    "geaendert": geaendert,
                    "ohne_inhalt": True,
                })
                continue

            dateiname_ohne_endung = os.path.splitext(os.path.basename(pfad))[0]

            for abschnitt in in_abschnitte_teilen(text):
                text_fuer_analyse = f"Dokument: {dateiname_ohne_endung}\nInhalt: {abschnitt}"
                neue_eintraege.append({
                    "datei": pfad,
                    "text": abschnitt,
                    "text_fuer_analyse": text_fuer_analyse,
                    "geaendert": geaendert,
                })

    # KI-Vektoren in kleinen Batches berechnen statt alles auf einmal.
    # Vorher wurde modell.encode() einmal für ALLE gesammelten Textabschnitte
    # aufgerufen - bei vielen Dateien konnte das (besonders im CPU-Modus,
    # siehe lade_modell()) mehrere Minuten dauern, OHNE dass die GUI
    # währenddessen einen Fortschritt anzeigen konnte. Es sah dann so aus,
    # als sei die App eingefroren. Jetzt wird in Batches gerechnet und nach
    # jedem Batch fortschritt_fn(None, ...) aufgerufen, damit die GUI eine
    # eigene, laufende Statusmeldung für diese Phase zeigen kann.
    # Zwischenspeicherung alle SPEICHER_INTERVALL Batches - ermöglicht
    # "Suchen während der Indexierung": statt bis zum kompletten Abschluss
    # (kann bei vielen Dateien über eine Stunde dauern) zu warten, kann man
    # schon nach den ersten fertig verarbeiteten Batches danach suchen,
    # während der Rest im Hintergrund weiterläuft.
    SPEICHER_INTERVALL = 3
    BATCH_GROESSE = 16
    zu_kodierende = [e for e in neue_eintraege if e.get("text_fuer_analyse")]
    if zu_kodierende:
        gesamt_batches = math.ceil(len(zu_kodierende) / BATCH_GROESSE)
        for batch_start in range(0, len(zu_kodierende), BATCH_GROESSE):
            batch = zu_kodierende[batch_start:batch_start + BATCH_GROESSE]
            texte = [e["text_fuer_analyse"] for e in batch]
            vektoren = modell.encode(texte, show_progress_bar=False, normalize_embeddings=True)
            for e, v in zip(batch, vektoren):
                e["vektor"] = v
                del e["text_fuer_analyse"]

            aktueller_batch = batch_start // BATCH_GROESSE + 1
            if fortschritt_fn:
                # idx=None signalisiert der GUI: das ist die Berechnungs-
                # Phase, nicht ein neu gelesenes Dokument - siehe gui.py.
                fortschritt_fn(None, f"Berechne KI-Vektoren (Batch {aktueller_batch}/{gesamt_batches})...")

            ist_letzter_batch = aktueller_batch == gesamt_batches
            if aktueller_batch % SPEICHER_INTERVALL == 0 or ist_letzter_batch:
                fertige_eintraege = [
                    e for e in neue_eintraege
                    if "vektor" in e or e.get("ohne_inhalt")
                ]
                zwischenstand = andere_eintraege + eintraege_dieser_ordner + fertige_eintraege
                speichere_index(zwischenstand)

    for e in neue_eintraege:
        e.pop("text_fuer_analyse", None)

    eintraege_dieser_ordner.extend(neue_eintraege)
    gesamt = andere_eintraege + eintraege_dieser_ordner
    speichere_index(gesamt)

    return gesamt


def datei_oeffnen(pfad):
    import platform
    try:
        system = platform.system()
        if system == "Darwin":
            subprocess.run(["open", pfad], check=True)
        elif system == "Windows":
            os.startfile(pfad)
        else:
            subprocess.run(["xdg-open", pfad], check=True)
    except Exception as e:
        print(f"Fehler beim Öffnen: {e}")


_modell_cache = None
_modell_lock = threading.Lock()


def geladenes_modell():
    """Lädt das Modell einmalig und cached es (Thread-sicher).

    Ohne den Lock könnten Suche und Indexierung, wenn sie gleichzeitig zum
    allerersten Mal starten, beide parallel lade_modell() aufrufen und das
    Modell doppelt laden (unnötiger Speicher-/Zeitverbrauch, im schlimmsten
    Fall doppelte Downloads beim ersten Start).
    """
    global _modell_cache
    if _modell_cache is None:
        with _modell_lock:
            if _modell_cache is None:  # Doppelt geprüft: evtl. hat ein
                # anderer Thread es inzwischen schon geladen, während wir
                # auf den Lock gewartet haben.
                _modell_cache = lade_modell()
    return _modell_cache


def anfrage_woerter(anfrage):
    rohe_woerter = re.findall(r"\w+", anfrage.lower())
    return [w for w in rohe_woerter if len(w) > 2 and w not in STOPWOERTER]


def _zeitraum_cutoff(zeitraum):
    """Wandelt eine Zeitraum-Auswahl der GUI in einen Unix-Timestamp um,
    ab dem eine Datei als 'im Zeitraum' gilt. None = kein Filter."""
    if not zeitraum or zeitraum == "Alle Zeiten":
        return None

    heute = datetime.date.today()
    if zeitraum == "7 Tage":
        return time.time() - 7 * 86400
    elif zeitraum == "Dieser Monat":
        start = datetime.date(heute.year, heute.month, 1)
        return time.mktime(start.timetuple())
    elif zeitraum == "Dieses Jahr":
        start = datetime.date(heute.year, 1, 1)
        return time.mktime(start.timetuple())
    return None


def suche_intern(anfrage, top_n=10, ausgeschlossene_typen=None, zeitraum=None):
    if not os.path.exists(INDEX_FILE):
        return None

    import numpy as np

    with open(INDEX_FILE, "rb") as f:
        eintraege = pickle.load(f)

    # Einträge ohne Vektor (z.B. Dateien, die beim Indexieren nicht gelesen
    # werden konnten) können nicht durchsucht werden - herausfiltern.
    eintraege = [e for e in eintraege if "vektor" in e]

    if ausgeschlossene_typen:
        eintraege = [
            e for e in eintraege
            if os.path.splitext(e["datei"])[1].lower() not in ausgeschlossene_typen
        ]

    cutoff = _zeitraum_cutoff(zeitraum)
    if cutoff is not None:
        eintraege = [e for e in eintraege if e.get("geaendert", 0) >= cutoff]

    if not eintraege:
        return []

    modell = geladenes_modell()
    anfrage_vektor = modell.encode([anfrage], normalize_embeddings=True)[0]
    such_woerter = anfrage_woerter(anfrage)

    rohe_treffer = []
    for e in eintraege:
        semantik_score = float(np.dot(anfrage_vektor, e["vektor"]))

        dateiname = os.path.basename(e["datei"]).lower()
        text_lc = e["text"].lower()

        kw_score = 0.0
        if such_woerter:
            treffer_anzahl = 0
            for w in such_woerter:
                if w in dateiname:
                    treffer_anzahl += 3.0
                elif w in text_lc:
                    treffer_anzahl += 1.5
            kw_score = (treffer_anzahl / len(such_woerter)) * 0.3

        gesamt_score = semantik_score + kw_score

        if gesamt_score > SCORE_SCHWELLE:
            rohe_treffer.append((gesamt_score, e))

    rohe_treffer.sort(key=lambda x: x[0], reverse=True)

    gesehene_dateien = set()
    eindeutige_ergebnisse = []

    for score, eintrag in rohe_treffer:
        datei_pfad = eintrag["datei"]
        if datei_pfad not in gesehene_dateien:
            gesehene_dateien.add(datei_pfad)
            eindeutige_ergebnisse.append((score, eintrag))

        if len(eindeutige_ergebnisse) >= top_n:
            break

    return eindeutige_ergebnisse


def aehnliche_dateien(pfad, top_n=10):
    """Findet Dateien, die INHALTLICH ähnlich zu 'pfad' sind - per
    Vektor-Ähnlichkeit statt nur Dateiname-Textsuche.

    Nutzt den bereits vorhandenen Embedding-Vektor der Datei aus dem Index
    (kein neuer API-Call/keine neue Berechnung nötig) und vergleicht ihn
    per Kosinus-Ähnlichkeit (Skalarprodukt, da alle Vektoren normalisiert
    sind) gegen alle anderen indexierten Dateien.

    Gibt eine Liste von (score, eintrag) zurück, im selben Format wie
    suche_intern(), damit die GUI dieselbe Ergebnisanzeige wiederverwenden
    kann. None, wenn die Datei nicht (mit Vektor) im Index ist.
    """
    if not os.path.exists(INDEX_FILE):
        return None

    import numpy as np

    pfad = os.path.abspath(os.path.expanduser(pfad))

    with open(INDEX_FILE, "rb") as f:
        eintraege = pickle.load(f)

    eintraege = [e for e in eintraege if "vektor" in e]

    eigene_vektoren = [e["vektor"] for e in eintraege if e["datei"] == pfad]
    if not eigene_vektoren:
        return None

    # Bei mehreren Textabschnitten derselben Datei: Durchschnittsvektor als
    # Repräsentation der ganzen Datei verwenden.
    referenz_vektor = np.mean(eigene_vektoren, axis=0)
    referenz_vektor = referenz_vektor / np.linalg.norm(referenz_vektor)

    rohe_treffer = []
    for e in eintraege:
        if e["datei"] == pfad:
            continue  # Datei nicht mit sich selbst vergleichen
        score = float(np.dot(referenz_vektor, e["vektor"]))
        rohe_treffer.append((score, e))

    rohe_treffer.sort(key=lambda x: x[0], reverse=True)

    gesehene_dateien = set()
    eindeutige_ergebnisse = []
    for score, eintrag in rohe_treffer:
        datei_pfad = eintrag["datei"]
        if datei_pfad not in gesehene_dateien:
            gesehene_dateien.add(datei_pfad)
            eindeutige_ergebnisse.append((score, eintrag))
        if len(eindeutige_ergebnisse) >= top_n:
            break

    return eindeutige_ergebnisse


def main():
    if len(sys.argv) < 2:
        return
    befehl = sys.argv[1]
    if befehl == "ordner" and len(sys.argv) >= 3 and sys.argv[2] == "hinzufuegen":
        befehl_ordner_hinzufuegen(sys.argv[3])


if __name__ == "__main__":
    main()